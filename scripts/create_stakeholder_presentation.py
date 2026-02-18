"""
Generate ERIS stakeholder PowerPoint presentation.
Run from project root:  python scripts/create_stakeholder_presentation.py
Output:  ERIS_Stakeholder_Presentation.pptx (in project root or output/)
"""

import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Install python-pptx:  pip install python-pptx")
    raise

PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_PATH = PROJECT_ROOT / "ERIS_Stakeholder_Presentation.pptx"

# ---------- Helpers ----------
def add_title_slide(prs, title, subtitle=""):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.add_paragraph() if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(14)
    return slide

def add_section_slide(prs, section_title):
    """Section divider slide."""
    layout = prs.slide_layouts[5]  # section header style if available
    try:
        slide = prs.slides.add_slide(layout)
    except Exception:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = section_title
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----- 1. Title -----
    add_title_slide(
        prs,
        "ERIS — Economic Regime Intelligence System",
        "Stakeholder demonstration · End-to-end pipeline, models & business value",
    )

    # ----- 2. Agenda -----
    add_content_slide(prs, "Agenda", [
        "Project introduction & business objectives",
        "Data sources & collection",
        "Feature engineering & preprocessing",
        "Models: Sentiment, Regime, Topics",
        "Stress level, AI briefing & dashboard",
        "Technical terminology (with business significance)",
        "KPI success factors & conclusion",
    ])

    # ----- 3. Project introduction -----
    add_content_slide(prs, "Project introduction", [
        "What: ERIS is an Economic Regime Intelligence System that turns financial text and market data into actionable risk and regime signals.",
        "Why: Stakeholders need a single view of market mood (Risk-On / Transitional / Risk-Off), stress level, and narrative themes—not raw news alone.",
        "Who: Built for investors, risk managers, and leadership who need explainable, data-driven insights on regime and sentiment.",
        "Outcome: A live dashboard (Streamlit) with stress gauge, regime over time, topic themes, and GPT-powered briefings.",
    ])

    # ----- 4. Business problem & objectives -----
    add_content_slide(prs, "Business problem & objectives", [
        "Problem: Dispersed data (news, Fed, earnings, markets) and no unified view of economic regime or narrative themes.",
        "Objective 1: Ingest multiple data sources and store them in a single pipeline.",
        "Objective 2: Derive sentiment and regime from text so stakeholders see mood and risk level.",
        "Objective 3: Surface key themes (topics) and link regime to market data for context.",
        "Objective 4: Deliver a stakeholder-ready dashboard with KPIs and AI briefings.",
    ])

    # ----- 5. Solution overview -----
    add_content_slide(prs, "Solution overview — high-level flow", [
        "Phase 1 — Collect: News (NewsAPI), Fed (scraper), Market (yfinance, FRED), Kaggle (financial news + earnings), optional CSV earnings.",
        "Phase 2 — Preprocess: Clean text, deduplicate, language filter, time-align → documents_processed.",
        "Phase 3 — NLP: FinBERT sentiment per document → nlp_signals; BERTopic for topic labels.",
        "Phase 4 — Regime: Daily sentiment features → HMM (3 states) → regime_states (Risk-On / Transitional / Risk-Off).",
        "Phase 5 — App: Stress score (0–100), regime chart, topics, market overlay, GPT risk briefing & mitigation paths.",
    ])

    # ----- 6. Data sources (detailed) -----
    add_content_slide(prs, "Data sources — where the data comes from", [
        "NewsAPI: Regime-relevant queries (inflation, recession, fed policy, credit markets, banking stress, yield curve, rate hike, QT, unemployment, earnings miss, default risk). Business significance: Real-time narrative that drives market sentiment.",
        "Federal Reserve: FOMC statements, speeches (fed_scraper). Business significance: Policy and forward guidance that move markets.",
        "Market (yfinance): SPY, VIX, GLD, TLT, HYG — OHLCV and derived returns. Business significance: Observable market state to compare with text-based regime.",
        "FRED: DGS10, DGS2, T10Y2Y, HY spread, UNRATE, CPI. Business significance: Macro context for stress and regime.",
        "Kaggle: Financial news datasets (e.g. sentiment-analysis-for-financial-news), earnings call transcripts. Business significance: Historical bulk for robust sentiment and topic modeling.",
        "Earnings: CSV (data/raw/earnings.csv) or Kaggle earnings. Business significance: Corporate narrative and forward-looking risk.",
    ])

    # ----- 7. Data pipeline overview -----
    add_content_slide(prs, "Data pipeline — order of operations", [
        "1. Schema: Create DB tables (raw_articles, fed_documents, earnings_transcripts, market_daily, documents_processed, nlp_signals, regime_states).",
        "2. Collectors: news_collector → fed_scraper → market_collector → kaggle_collector → earnings (CSV).",
        "3. Preprocess: Run on raw_articles, fed_documents, earnings_transcripts → documents_processed.",
        "4. Sentiment: FinBERT on documents_processed → nlp_signals (per-doc then aggregated by date).",
        "5. Regime: Daily features from nlp_signals → HMM → regime_states.",
        "6. Topics: BERTopic on documents_processed → topic_hint per document.",
        "7. Dashboard: Streamlit reads from DB; stress and GPT briefing use regime + sentiment + topics.",
    ])

    # ----- 8. Feature engineering — preprocessing -----
    add_content_slide(prs, "Feature engineering — preprocessing (after collection)", [
        "Cleaning: Strip HTML (trafilatura), remove boilerplate, normalize whitespace; sentence segmentation (spaCy when available). Output: content_clean, content_sentences, word_count.",
        "Deduplication: MinHash fingerprint + Jaccard threshold (e.g. 0.85) so near-duplicate articles are dropped. Business significance: Reduces noise and double-counting of the same story.",
        "Language filter: Keep English-only (langdetect). Minimum word count (e.g. 20) to drop snippets.",
        "Time alignment: Map published timestamps to calendar date (align_publish_to_date) for consistent daily aggregation. Business significance: Correct date ensures regime and sentiment align to the right trading day.",
        "Output table: documents_processed (source_id, source_table, source_type, title, content_clean, content_sentences, published_date, topic_hint, minhash_fingerprint).",
    ])

    # ----- 9. Feature engineering — market -----
    add_content_slide(prs, "Feature engineering — market data", [
        "OHLCV: Open, high, low, close, volume from yfinance (and FRED for macro series).",
        "Returns: returns_1d, returns_5d, returns_21d, returns_63d (percent change over 1, 5, 21, 63 days). Business significance: Momentum and horizon-specific performance.",
        "Realized volatility: 21-day rolling standard deviation of daily returns (realized_vol_21d). Business significance: Captures recent market uncertainty for stress context.",
        "Stored in: market_daily (per ticker/date), macro_indicators (FRED series).",
    ])

    # ----- 10. Model 1 — Sentiment (FinBERT) -----
    add_content_slide(prs, "Model 1 — Sentiment (FinBERT)", [
        "Technical: FinBERT (ProsusAI/finbert) is a BERT model fine-tuned on financial text. Input: sentence(s) from document. Output: (score, p_positive, p_negative, p_neutral); score = p_pos − p_neg on scale −1 to +1.",
        "Granular: We split each document into sentences, score each (up to 50 per doc), then average. One row per document in nlp_signals with date, source_type, sentiment_score, and probabilities.",
        "Business significance: Converts unstructured news/Fed/earnings text into a numeric sentiment signal. Negative sentiment clusters often precede or coincide with risk-off regimes; positive with risk-on.",
    ])

    # ----- 11. Model 2 — Regime (HMM) -----
    add_content_slide(prs, "Model 2 — Regime detection (HMM)", [
        "Technical: Daily features from nlp_signals: sentiment_mean, sentiment_std, sentiment_drift (21-day rolling mean diff). Features are z-scored. Gaussian HMM (3 components, full covariance) from hmmlearn; states are labeled Risk-On, Transitional, Risk-Off using emission means.",
        "Granular: One regime state per calendar day; probability of Risk-Off (and confidence) stored in regime_states. Optional: ruptures for change-point detection; configurable n_iter (e.g. 100).",
        "Business significance: Gives a single, interpretable regime label per day so stakeholders can see whether the economy is in risk-on, transitional, or risk-off mode—and how confident the model is.",
    ])

    # ----- 12. Model 3 — Topics (BERTopic) -----
    add_content_slide(prs, "Model 3 — Topics (BERTopic)", [
        "Technical: Sentence-transformers (e.g. all-MiniLM-L6-v2) for embeddings; UMAP for dimensionality reduction; HDBSCAN for clustering (no fixed number of topics). Each document gets a topic_id; BERTopic generates a label per topic. We store topic_hint on documents_processed.",
        "Granular: Adaptive min_cluster_size / min_samples for small corpora so more documents get real themes instead of one big outlier bucket. Humanized labels (e.g. Bitcoin & crypto, Federal Reserve) for the dashboard.",
        "Business significance: Surfaces dominant narrative themes (inflation, Fed, credit, etc.) so stakeholders see what the corpus is about and how it ties to regime and stress.",
    ])

    # ----- 13. Stress level & KPIs -----
    add_content_slide(prs, "Stress level & KPI dashboard", [
        "Stress score (0–100): Derived from current regime Risk-Off probability and latest daily sentiment. Scale: 0–25 Low, 25–50 Elevated, 50–75 High, 75–100 Critical. Business significance: Single number for how stressed the market is; drives suggested actions (e.g. capital preservation at Critical).",
        "KPI page: Raw articles, processed docs, regime days, NLP signals, topic diversity vs targets (e.g. ≥1000 raw, ≥500 processed, ≥90 regime days). Gauges show % of target achieved.",
        "Success summary: All benchmark targets met → pipeline is industry-ready; below benchmark → run full pipeline (e.g. Kaggle + topic labels) for production quality.",
    ])

    # ----- 14. LLM / AI briefing -----
    add_content_slide(prs, "AI briefing & mitigation paths", [
        "Input: Current regime, regime trend, sentiment trend, topic summary, stress info. Sent to GPT-4 (configurable model) with a structured prompt.",
        "Output: scenario_summary, risk_briefing (full narrative), early_warnings, precautions, mitigation_paths. Cached per date in llm_briefings.",
        "Business significance: Converts technical signals into plain-language risk briefing and concrete mitigation steps for leadership and risk committees.",
    ])

    # ----- 15. End-to-end flow -----
    add_content_slide(prs, "End-to-end flow (summary)", [
        "Data → Preprocess → Sentiment → Regime (+ Topics, Market) → Stress + Briefing → Dashboard.",
        "Stakeholders see: stress gauge, regime over time, sentiment trend, topic distribution, market vs regime overlay, AI risk briefing and mitigation paths, and KPI success page.",
        "Pipeline can be run locally (run_all_data.py) or from the app (Run pipeline now) with configurable limits and optional topic labels.",
    ])

    # ----- 16. Technical terminology — glossary -----
    add_content_slide(prs, "Technical terminology — key terms", [
        "FinBERT: Pre-trained BERT for financial sentiment. Business: Enables consistent, comparable sentiment from news/Fed/earnings.",
        "HMM (Hidden Markov Model): Probabilistic model over a sequence of states. Business: Regime is a sequence of daily states (Risk-On / Transitional / Risk-Off).",
        "BERTopic: Topic model using embeddings + UMAP + HDBSCAN. Business: Automatic theme discovery without pre-defined categories.",
        "MinHash: Sketch for fast similarity (Jaccard). Business: Deduplication so the same story from multiple sources counts once.",
        "NLP signals: Document-level sentiment and metadata stored by date. Business: Input for daily regime and dashboards.",
        "Regime state: Single daily label (Risk-On / Transitional / Risk-Off) with probability. Business: Clear, communicable market mood.",
    ])

    # ----- 17. Technical terminology (continued) -----
    add_content_slide(prs, "Technical terminology (continued)", [
        "Stress score: 0–100 index from regime probability and sentiment. Business: One number for risk committees and alerts.",
        "documents_processed: Cleaned, deduplicated, time-aligned text ready for NLP. Business: Quality input for models.",
        "realized_vol_21d: 21-day rolling volatility of returns. Business: Captures recent market uncertainty.",
        "sentiment_drift: Change in 21-day rolling mean sentiment. Business: Captures shift in narrative tone for regime transitions.",
    ])

    # ----- 18. KPI success factors -----
    add_content_slide(prs, "KPI success factors", [
        "Data coverage: Raw articles ≥1000, processed docs ≥500, regime days ≥90, NLP signals ≥500, topic themes ≥3.",
        "Date range: Documents and regime series should span the period of interest (e.g. 2025–2026) for continuous insight.",
        "Topic diversity: Multiple distinct themes (not everything in Other) so narrative analysis is meaningful.",
        "Regime continuity: Many days of regime_states so the chart shows trends, not sparse spikes.",
        "Benchmark: When all targets are met, the project is considered successful and industry-relevant for stakeholder use.",
    ])

    # ----- 19. Conclusion -----
    add_content_slide(prs, "Conclusion — industry relevance", [
        "ERIS turns multi-source financial text and market data into a single, explainable view of economic regime, stress, and themes.",
        "Stakeholders get: stress level, regime over time, topic distribution, market linkage, and AI-generated risk briefing with mitigation paths.",
        "Success is measured by data coverage and KPI benchmarks; the dashboard and pipeline are designed for production use (local or Streamlit Cloud).",
        "Next steps: Maintain pipeline runs (collectors + preprocess + sentiment + regime + topics), monitor KPIs, and iterate on benchmarks for your organization.",
    ])

    # ----- 20. Thank you -----
    add_title_slide(prs, "Thank you", "Questions & discussion")

    # ----- Save -----
    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")
    return OUTPUT_PATH


if __name__ == "__main__":
    main()
